"""
==========================================================
Document Benchmark
Classic Parser Extractor
==========================================================

Extract text from different document formats using
specialized Python libraries.

Supported formats:
- PDF  -> PyMuPDF
- DOCX -> python-docx
- PPTX -> python-pptx
- TXT  -> Python built-in
"""

from pathlib import Path
import time

import fitz
from docx import Document
from pptx import Presentation

def save_output(file_path, text, extension="txt"):
    """
    Save extracted text into the output folder.
    """

    output_dir = Path("../output/parser")
    output_dir.mkdir(parents=True, exist_ok=True)

    output_file = output_dir / f"{Path(file_path).stem}.{extension}"

    with open(output_file, "w", encoding="utf-8") as f:
        f.write(text)

# ==========================================================
# PDF
# ==========================================================

def extract_pdf(file_path):

    start = time.time()

    document = fitz.open(file_path)

    text = ""
    page_count = len(document)

    for page in document:
        text += page.get_text() + "\n"

    document.close()
    save_output(file_path, text.strip())
    return {
        "tool": "Classic Parser",
        "file_name": Path(file_path).name,
        "file_type": "pdf",
        "processing_time": round(time.time() - start, 3),
        "page_count": page_count,
        "text": text.strip()
    }


# ==========================================================
# DOCX
# ==========================================================

def extract_docx(file_path):

    start = time.time()

    document = Document(file_path)

    paragraphs = []

    for p in document.paragraphs:
        paragraphs.append(p.text)

    text = "\n".join(paragraphs)
    save_output(file_path, text)

    return {
        "tool": "Classic Parser",
        "file_name": Path(file_path).name,
        "file_type": "docx",
        "processing_time": round(time.time() - start, 3),
        "page_count": None,
        "text": text.strip()
    }


# ==========================================================
# PPTX
# ==========================================================

def extract_pptx(file_path):

    start = time.time()

    presentation = Presentation(file_path)

    slides_text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                slides_text.append(shape.text)

    text = "\n".join(slides_text)
    save_output(file_path, text)

    return {
        "tool": "Classic Parser",
        "file_name": Path(file_path).name,
        "file_type": "pptx",
        "processing_time": round(time.time() - start, 3),
        "page_count": len(presentation.slides),
        "text": text.strip()
    }


# ==========================================================
# TXT
# ==========================================================

def extract_txt(file_path):

    start = time.time()

    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()
    save_output(file_path, text)

    return {
        "tool": "Classic Parser",
        "file_name": Path(file_path).name,
        "file_type": "txt",
        "processing_time": round(time.time() - start, 3),
        "page_count": None,
        "text": text.strip()
    }


# ==========================================================
# Main extractor
# ==========================================================

def extract_document(file_path):

    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return extract_pdf(file_path)

    elif extension == ".docx":
        return extract_docx(file_path)

    elif extension == ".pptx":
        return extract_pptx(file_path)

    elif extension == ".txt":
        return extract_txt(file_path)

    else:
        raise ValueError(f"Unsupported format: {extension}")

